# app.py
from pptx import Presentation
from pptx.util import Inches, Pt
import streamlit as st
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import io
import tempfile
import os
import seaborn as sns
from matplotlib.figure import Figure
from pandas.api.types import is_numeric_dtype, is_string_dtype

st.set_page_config(layout="wide", page_title="CSVisi0n - Quick Data Analyzer")

st.title("CSVisi0n — CSV → Excel + Visualizations + AI Summary")
st.caption("Upload CSV, inspect, visualize, get a textual summary and export an Excel report. Deadline-mode ready.")

uploaded = st.file_uploader("Upload CSV file", type=["csv", "txt"])
if not uploaded:
    st.info("Upload a CSV to start. If your CSV is large, try a small sample for the first run.")
    st.stop()

@st.cache_data
def load_csv(file):
    try:
        return pd.read_csv(file)
    except Exception:
        file.seek(0)
        return pd.read_csv(file, encoding='latin1', error_bad_lines=False)

df = load_csv(uploaded)
st.success(f"Loaded {len(df)} rows and {len(df.columns)} columns")

# Basic preview and controls
with st.expander("Data preview & basic info", expanded=True):
    st.subheader("Preview")
    st.dataframe(df.head(100))
    c1, c2, c3 = st.columns([2,1,1])
    with c1:
        st.write("Columns:", list(df.columns))
    with c2:
        st.write("Shape:", df.shape)
    with c3:
        st.write("Memory:", f"{df.memory_usage(deep=True).sum()/1024**2:.2f} MB")

# Summary stats
with st.expander("Automatic summary (counts, dtypes, missing)", expanded=False):
    dtypes = df.dtypes.astype(str)
    missing = df.isna().sum()
    summary_tbl = pd.DataFrame({"dtype": dtypes, "missing": missing, "missing_pct": missing/len(df)*100})
    st.dataframe(summary_tbl.sort_values("missing", ascending=False))

# Helper plotting functions
def plot_numeric_hist(col: str):
    fig = Figure(figsize=(4,3))
    ax = fig.subplots()
    sns.histplot(df[col].dropna(), ax=ax, kde=True)
    ax.set_title(f"Histogram — {col}")
    return fig

def plot_categorical_bar(col: str, top_n=20):
    vc = df[col].value_counts(dropna=False).head(top_n)
    fig = Figure(figsize=(4,3))
    ax = fig.subplots()
    vc.plot.bar(ax=ax)
    ax.set_title(f"Top values — {col}")
    return fig

def plot_corr_heatmap(numeric_df: pd.DataFrame):
    corr = numeric_df.corr()
    fig = Figure(figsize=(8,6))
    ax = fig.subplots()
    sns.heatmap(corr, annot=True, fmt=".2f", ax=ax, cmap="vlag", center=0)
    ax.set_title("Correlation heatmap")
    return fig, corr

def plot_missing_map():
    fig = Figure(figsize=(8,2))
    ax = fig.subplots()
    sns.heatmap(df.isna().T, cbar=False, ax=ax)
    ax.set_ylabel("Columns")
    ax.set_title("Missing values map")
    return fig

# Visualize: pick top numeric and categorical columns
numeric_cols = [c for c in df.columns if is_numeric_dtype(df[c])]
cat_cols = [c for c in df.columns if (is_string_dtype(df[c]) or df[c].nunique() < 30)]

st.subheader("Recommended Visualizations")
cols = st.columns(2)
with cols[0]:
    st.write("Numeric columns")
    if numeric_cols:
        for col in numeric_cols[:6]:
            st.pyplot(plot_numeric_hist(col))
    else:
        st.info("No numeric columns detected.")
with cols[1]:
    st.write("Categorical columns")
    if cat_cols:
        for col in cat_cols[:6]:
            st.pyplot(plot_categorical_bar(col))
    else:
        st.info("No categorical columns detected.")

# Correlation & missingness
if numeric_cols:
    corr_fig, corr_df = plot_corr_heatmap(df[numeric_cols])
    st.pyplot(corr_fig)
else:
    st.info("No numeric columns available for correlation heatmap.")

st.pyplot(plot_missing_map())

# Automated "AI" summary generator (deterministic)
def generate_summary(df: pd.DataFrame, top_n_corr=5):
    nrows, ncols = df.shape
    out = []
    out.append(f"The dataset has {nrows} rows and {ncols} columns.")
    # missingness
    miss = df.isna().sum()
    total_missing = miss.sum()
    out.append(f"Total missing values: {int(total_missing)} ({total_missing / (nrows*ncols) * 100:.2f}% of cells).")
    high_missing = miss[miss > 0.3 * nrows].sort_values(ascending=False)
    if len(high_missing):
        out.append(f"Columns with >30% missing: {list(high_missing.index[:10])}")
    # types
    dtypes = df.dtypes.value_counts().to_dict()
    out.append("Column types: " + ", ".join([f"{k}:{v}" for k,v in dtypes.items()]))
    # numeric insights
    if len(numeric_cols):
        out.append(f"There are {len(numeric_cols)} numeric columns. Quick stats (mean, std, min, median, max) for a few:")
        stats = df[numeric_cols].describe().T
        sample_stats = stats.iloc[:min(5, len(stats))]
        for idx, row in sample_stats.iterrows():
            out.append(f"- {idx}: mean={row['mean']:.3g}, std={row['std']:.3g}, min={row['min']:.3g}, median={row['50%']:.3g}, max={row['max']:.3g}")
        # outliers via IQR
        outlier_cols = []
        for c in numeric_cols:
            series = df[c].dropna()
            if series.empty: continue
            q1, q3 = np.percentile(series, [25,75])
            iqr = q3 - q1
            if iqr == 0: continue
            lower = q1 - 1.5*iqr
            upper = q3 + 1.5*iqr
            outlier_frac = ((series < lower) | (series > upper)).mean()
            if outlier_frac > 0.01:
                outlier_cols.append((c, outlier_frac))
        if outlier_cols:
            out.append("Columns with >1% outliers (IQR rule): " + ", ".join([f"{c}({f:.2%})" for c,f in outlier_cols[:10]]))
    # correlation
    if len(numeric_cols) > 1:
        corr = df[numeric_cols].corr().abs().unstack().sort_values(ascending=False)
        corr = corr[corr < 1].drop_duplicates()
        top = corr.head(top_n_corr)
        if not top.empty:
            out.append("Top absolute correlations:")
            for (a,b),val in top.items():
                out.append(f"- {a} vs {b}: {val:.2f}")
    # categorical tips
    if cat_cols:
        out.append(f"{len(cat_cols)} categorical-ish columns detected (low-cardinality or string). Consider bar counts, stacked bars, or pivot tables.")
    # action items
    out.append("Suggested next steps: handle columns with large missingness (drop/fill), ensure date columns parsed, consider aggregations for high-cardinality categories, create targeted visualizations for domain-specific questions.")
    return "\n".join(out)

summary_text = generate_summary(df)
st.subheader("AI Summary (automated)")
st.text_area("Summary", value=summary_text, height=220)
# Prepare images for PPT
img_paths = []

# numeric histograms
for col in numeric_cols[:6]:
    fig = plot_numeric_hist(col)
    p = f"hist_{col}.png"
    fig.savefig(p, bbox_inches='tight')
    img_paths.append(p)

# categorical charts
for col in cat_cols[:6]:
    fig = plot_categorical_bar(col)
    p = f"cat_{col}.png"
    fig.savefig(p, bbox_inches='tight')
    img_paths.append(p)

# correlation heatmap
if numeric_cols:
    corr_fig, _ = plot_corr_heatmap(df[numeric_cols])
    p = "corr.png"
    corr_fig.savefig(p, bbox_inches='tight')
    img_paths.append(p)

# missing map
p = "missing.png"
plot_missing_map().savefig(p, bbox_inches='tight')
img_paths.append(p)

# Create Excel report with sheets and embedded images
def create_excel_report(df: pd.DataFrame):
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".xlsx")
    path = tmp.name
    writer = pd.ExcelWriter(path, engine='xlsxwriter')
    df.to_excel(writer, sheet_name='raw_data', index=False)
    # write summary sheet
    workbook = writer.book
    summary_sheet = workbook.add_worksheet("summary")
    writer.sheets['summary'] = summary_sheet
    summary_sheet.write(0,0,"AI Summary")
    for i, line in enumerate(generate_summary(df).splitlines(), start=1):
        summary_sheet.write(i,0,line)
    # generate and save some plots as images and insert
    img_dir = tempfile.mkdtemp()
    img_paths = []
    # numeric histograms
    for i, col in enumerate(numeric_cols[:6]):
        fig = plot_numeric_hist(col)
        p = os.path.join(img_dir, f"hist_{i}.png")
        fig.savefig(p, bbox_inches='tight')
        img_paths.append(p)
    for i, col in enumerate(cat_cols[:6]):
        fig = plot_categorical_bar(col)
        p = os.path.join(img_dir, f"cat_{i}.png")
        fig.savefig(p, bbox_inches='tight')
        img_paths.append(p)
    # correlation
    if numeric_cols:
        corr_fig, _ = plot_corr_heatmap(df[numeric_cols])
        p = os.path.join(img_dir, "corr.png")
        corr_fig.savefig(p, bbox_inches='tight')
        img_paths.append(p)
    missing_path = os.path.join(img_dir, "missing.png")
    plot_missing_map().savefig(missing_path, bbox_inches='tight')
    img_paths.append(missing_path)
    # insert images into summary sheet
    r = 0
    for p in img_paths:
        try:
            summary_sheet.insert_image(r, 2, p, {'x_scale': 0.6, 'y_scale': 0.6})
            r += 15
        except Exception:
            pass
    writer.close()
    return path
    from pptx import Presentation
from pptx.util import Inches, Pt

def create_ppt_report(df, summary_text, image_paths):
    prs = Presentation()
    
    # TITLE SLIDE
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "CSVisi0n – Data Analysis Report"
    subtitle.text = "Auto-generated PowerPoint Summary"

    # SUMMARY SLIDE
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    body = slide.shapes.placeholders[1]
    title.text = "AI Summary"
    tf = body.text_frame
    for line in summary_text.splitlines():
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(14)

    # VISUALIZATION SLIDES
    for img in image_paths:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = img.split("/")[-1]
        slide.shapes.add_picture(img, Inches(1), Inches(1.3), width=Inches(8))

    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    prs.save(tmp.name)
    return tmp.name

if st.button("Generate Report & Prepare Excel"):
    with st.spinner("Creating Excel report..."):
        excel_path = create_excel_report(df)
    st.success("Excel report created.")
    with open(excel_path, "rb") as f:
        data = f.read()
    b64 = io.BytesIO(data)
    st.download_button("Download Excel report (.xlsx)", data=b64, file_name="csvision_report.xlsx", mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")
if st.button("Generate PowerPoint (.pptx)"):
    with st.spinner("Creating PPT..."):
        ppt_path = create_ppt_report(df, summary_text, img_paths)
    st.success("PPT ready!")
    with open(ppt_path, "rb") as f:
        st.download_button("Download PowerPoint", f, file_name="csvision_report.pptx")

st.markdown("---")
st.caption("If you want, I can generate a pared-down README and a 1-minute demo script you can use in the presentation. Want them?")
